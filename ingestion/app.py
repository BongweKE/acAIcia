import os
import json
import logging
from typing import List
import modal

logger = logging.getLogger(__name__)

image = modal.Image.debian_slim().pip_install(
    "torch",
    "sentence-transformers",
    "PyMuPDF",
    "python-docx",
    "python-pptx",
    "langchain-text-splitters",
    "supabase"
)

app = modal.App("acaicia-ingestion")
vol = modal.Volume.from_name("acaicia-data-volume", create_if_missing=True)

secrets = [
    modal.Secret.from_name("acaicia-db-secrets")
]

@app.function(image=image, gpu="T4", volumes={"/data": vol}, secrets=secrets, timeout=3600)
def process_documents():
    import glob
    import fitz
    from docx import Document
    from pptx import Presentation
    from supabase import create_client, Client
    from sentence_transformers import SentenceTransformer
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    
    SUPABASE_URL = os.environ.get("SUPABASE_URL")
    SUPABASE_KEY = os.environ.get("SUPABASE_KEY")
    supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
    
    logging.info("Loading embedding model BAAI/bge-base-en-v1.5 onto T4 GPU...")
    model = SentenceTransformer('BAAI/bge-base-en-v1.5')
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=2500,
        chunk_overlap=250,
        length_function=len
    )
    
    # Reload Volume cleanly before asserting states
    vol.reload()
    
    state_file = "/data/ingestion_state.json"
    state = {}
    if os.path.exists(state_file):
        with open(state_file, 'r') as f:
            try:
                state = json.load(f)
            except json.JSONDecodeError:
                state = {}
                
    metadata_file = "/data/metadata.json"
    catalog_meta = {}
    if os.path.exists(metadata_file):
        with open(metadata_file, 'r') as f:
            try:
                catalog_meta = json.load(f)
            except json.JSONDecodeError:
                logging.warning(f"Failed to parse {metadata_file}. Using default metadata mappings.")
            
    files = []
    for ext in ['*.pdf', '*.md', '*.docx', '*.pptx']:
        files.extend(glob.glob(f"/data/{ext}"))
        
    if not files:
        logging.info("No files found in /data volume.")
        return
        
    for file_path in files:
        filename = os.path.basename(file_path)
        if state.get(filename, {}).get("status") == "Success":
            logging.info(f"Skipping already processed file: {filename}")
            continue
            
        logging.info(f"Processing: {filename}")
        title, ext = os.path.splitext(filename)
        ext = ext.lower()
        text = ""
        
        try:
            if ext == '.pdf':
                doc = fitz.open(file_path)
                text = "".join(page.get_text() + "\n" for page in doc)
            elif ext == '.md':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            elif ext == '.docx':
                doc = Document(file_path)
                text = "".join(p.text + "\n" for p in doc.paragraphs)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                            
            if not text.strip():
                raise ValueError("No text extracted")
                
            chunks = text_splitter.split_text(text)
            num_chunks = len(chunks)
            logging.info(f"Generated {num_chunks} chunks. Encoding vectors natively...")
            
            # Embed all chunks
            embeddings = model.encode(chunks, convert_to_numpy=True)
            
            # Insert catalog metadata
            file_meta = catalog_meta.get(filename, {})
            doc_data = {
                "title": file_meta.get("title", title.replace('_', ' ').title()),
                "authors": file_meta.get("authors", []),
                "publication_year": file_meta.get("publication_year", None),
                "topic_keywords": file_meta.get("topic_keywords", []),
                "url_link": file_meta.get("url_link", None),
                "doi": file_meta.get("doi", None)
            }
            res = supabase.table("documents_catalog").insert(doc_data).execute()
            doc_id = res.data[0]['id']
            
            # Insert chunks payload sequentially
            for i, chunk in enumerate(chunks):
                chunk_data = {
                    "document_id": doc_id,
                    "chunk_text": chunk,
                    "embedding": embeddings[i].tolist()
                }
                supabase.table("document_embeddings").insert(chunk_data).execute()
                
            # Document logs update
            supabase.table("ingestion_logs").insert({
                "filename": filename,
                "chunks_created": num_chunks,
                "status": "Success",
                "error_message": None
            }).execute()
            
            state[filename] = {"status": "Success"}
            logging.info(f"Successfully finished embedding {filename}")
            
        except Exception as e:
            logging.error(f"Failed {filename}: {e}")
            supabase.table("ingestion_logs").insert({
                "filename": filename,
                "chunks_created": 0,
                "status": "Failed",
                "error_message": str(e)
            }).execute()
            state[filename] = {"status": "Failed", "error": str(e)}
            
    with open(state_file, 'w') as f:
        json.dump(state, f)
        
    vol.commit()
    logging.info("Volume state committed and process fully concluded.")

@app.local_entrypoint()
def main():
    print("Triggering remote document processing cluster directly...")
    process_documents.remote()
